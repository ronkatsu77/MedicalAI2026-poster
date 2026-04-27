import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from lxml import etree

DML = 'http://schemas.openxmlformats.org/drawingml/2006/main'
PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
PIC = 'http://schemas.openxmlformats.org/drawingml/2006/picture'
OOXML = DML  # alias for fill XML helpers

# ── Pink ribbon palette ────────────────────────────────────────────────────
# Section header bars:  deep rose → soft pink (left→right gradient)
HDR_START  = 'B5396A'   # deep dusty rose
HDR_END    = 'E090B8'   # soft petal pink

# まとめ background:  medium rose → pale pink
MTM_START  = 'C26090'
MTM_END    = 'E8A5C0'

# Title card full background:  dark maroon → vibrant rose (left→right)
TITLE_START = '8B1A42'
TITLE_END   = 'D96B9E'

# Accent line / affiliation text (was LT_BLUE D6E4F0)
ACCENT_PINK = 'FBD5E9'   # very light petal pink


# ── XML helpers ────────────────────────────────────────────────────────────
def grad(c1, c2, ang=0):
    """Linear gradient fill.  ang=0 → left→right."""
    return etree.fromstring(f'''<a:gradFill xmlns:a="{OOXML}" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>
  </a:gsLst>
  <a:lin ang="{ang}" scaled="0"/>
</a:gradFill>''')

def solid(c):
    return etree.fromstring(
        f'<a:solidFill xmlns:a="{OOXML}"><a:srgbClr val="{c}"/></a:solidFill>')

def _remove_fills(container):
    for tag in ['solidFill', 'gradFill', 'pattFill', 'blipFill', 'noFill']:
        for el in list(container.findall(f'{{{OOXML}}}{tag}')):
            container.remove(el)

def _insert_fill(container, fill_el, after_tag=None):
    """Insert fill_el after after_tag child (e.g. xfrm), or at index 0."""
    ref = container.find(f'{{{OOXML}}}{after_tag}') if after_tag else None
    pos = (list(container).index(ref) + 1) if ref is not None else 0
    container.insert(pos, fill_el)

def set_fill_spPr(spPr, fill_el):
    _remove_fills(spPr)
    _insert_fill(spPr, fill_el, after_tag='xfrm')

def set_fill_bgPr(bgPr, fill_el):
    _remove_fills(bgPr)
    bgPr.insert(0, fill_el)


def get_spPr(sp):
    # p:spPr (shapes), p:spPr in PIC namespace (pictures), or fallback via xpath
    for ns in [PML, PIC, DML]:
        found = sp.find(f'{{{ns}}}spPr')
        if found is not None:
            return found
    # lxml fallback: any element named spPr
    hits = sp.xpath('.//*[local-name()="spPr"]')
    return hits[0] if hits else None


def recolor_slide_shapes(slide, color_map):
    """Replace solid fills in shapes according to {old_HEX_upper: new_fill_element}."""
    for s in slide.shapes:
        spPr = get_spPr(s._element)
        if spPr is None:
            continue
        sf = spPr.find(f'{{{OOXML}}}solidFill')
        if sf is None:
            continue
        srgb = sf.find(f'{{{OOXML}}}srgbClr')
        if srgb is None:
            continue
        val = srgb.get('val', '').upper()
        if val in color_map:
            set_fill_spPr(spPr, color_map[val])
            print(f'  Recolored shape: {val} → {s.text.strip()[:30] if hasattr(s,"text") else ""}')


def recolor_text(slide, old_hex, new_rgb):
    """Replace explicit text run colors."""
    old = old_hex.upper()
    for s in slide.shapes:
        if not hasattr(s, 'text_frame'):
            continue
        for para in s.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if str(run.font.color.rgb).upper() == old:
                        run.font.color.rgb = new_rgb
                except Exception:
                    pass


# ════════════════════════════════════════════════════════════════════════════
# 1.  MAIN POSTER  3poster_v4.pptx → 3poster_pink.pptx
# ════════════════════════════════════════════════════════════════════════════
print('── Main poster ─────────────────────────────')
prs = Presentation('3poster_v4.pptx')
slide = prs.slides[0]

# Slide background: barely-blush (almost white)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF8)

# Recolor shapes
recolor_slide_shapes(slide, {
    '1F3864': grad(HDR_START, HDR_END),   # navy section headers → rose gradient
    '2E74B5': grad(MTM_START, MTM_END),   # blue まとめ bg       → soft rose gradient
})

prs.save('3poster_pink.pptx')
print('Saved: 3poster_pink.pptx')


# ════════════════════════════════════════════════════════════════════════════
# 2.  TITLE CARD  3poster_title_70x20.pptx → 3poster_title_pink.pptx
# ════════════════════════════════════════════════════════════════════════════
print()
print('── Title card ──────────────────────────────')
prs_t = Presentation('3poster_title_70x20.pptx')
slide_t = prs_t.slides[0]

# Slide background: deep maroon → vibrant rose gradient (left→right)
# background._element = <p:cSld>, which contains <p:bg><p:bgPr>
cSld  = slide_t.background._element
bg_el = cSld.find(f'{{{PML}}}bg')
bgPr  = bg_el.find(f'{{{PML}}}bgPr') if bg_el is not None else None
if bgPr is not None:
    set_fill_bgPr(bgPr, grad(TITLE_START, TITLE_END, ang=0))
    print('  Title card background → gradient')
else:
    slide_t.background.fill.solid()
    slide_t.background.fill.fore_color.rgb = RGBColor(0xB5, 0x39, 0x6A)
    print('  Title card background → solid deep rose (bgPr not found)')

# Accent line + other shapes that are navy
recolor_slide_shapes(slide_t, {
    '1F3864': grad(TITLE_START, TITLE_END),   # any navy rect
    'D6E4F0': solid(ACCENT_PINK),              # lt-blue accent line → light petal pink
})

# Affiliation text color: LT_BLUE → petal pink
recolor_text(slide_t, 'D6E4F0', RGBColor(0xFB, 0xD5, 0xE9))

prs_t.save('3poster_title_pink.pptx')
print('Saved: 3poster_title_pink.pptx')
