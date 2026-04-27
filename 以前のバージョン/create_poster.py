import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

IMGS = 'extracted_images/'

NAVY     = RGBColor(0x1F, 0x38, 0x64)
MED_BLUE = RGBColor(0x2E, 0x74, 0xB5)
LT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1A, 0x1A, 0x1A)
LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
BORDER   = RGBColor(0xAD, 0xC8, 0xE0)
FONT     = 'メイリオ'

# ── Primitive helpers ──────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill, border=None, bw=0.75):
    s = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s

def add_tb(slide, l, t, w, h, text, fs=11, bold=False,
           color=None, align=PP_ALIGN.LEFT):
    color = color or DARK
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left  = Cm(0.2)
    tf.margin_top   = Cm(0.1)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(fs)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = FONT
    return tb

def add_lines(slide, l, t, w, h, lines, bg=None, border=None):
    """lines: list of (text, bold, fontsize) tuples or plain str"""
    if bg:
        s = add_rect(slide, l, t, w, h, bg, border, 0.5)
        tf = s.text_frame
    else:
        tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
        tf = tb.text_frame
    tf.word_wrap    = True
    tf.margin_left  = Cm(0.25)
    tf.margin_top   = Cm(0.15)
    tf.margin_bottom= Cm(0.1)
    first = True
    for item in lines:
        text = item if isinstance(item, str) else item[0]
        bold = False  if isinstance(item, str) else item[1]
        fs   = 10.5  if isinstance(item, str) else (item[2] if len(item)>2 else 10.5)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(fs)
        run.font.bold  = bold
        run.font.color.rgb = DARK
        run.font.name  = FONT

def sec_hdr(slide, l, t, w, text, fs=14, bg=None, fg=WHITE, h=0.9):
    bg = bg or NAVY
    s = add_rect(slide, l, t, w, h, bg)
    tf = s.text_frame
    tf.margin_left  = Cm(0.35)
    tf.margin_top   = Cm(0.05)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(fs)
    run.font.bold  = True
    run.font.color.rgb = fg
    run.font.name  = FONT

def fig(slide, path, l, t, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(path, Cm(l), Cm(t), Cm(w), Cm(h))
    elif w:
        return slide.shapes.add_picture(path, Cm(l), Cm(t), width=Cm(w))
    elif h:
        return slide.shapes.add_picture(path, Cm(l), Cm(t), height=Cm(h))
    return slide.shapes.add_picture(path, Cm(l), Cm(t))

# ── Build poster ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Cm(45)
prs.slide_height = Cm(95)

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

# ================================================================
# HEADER  0 – 12.9 cm
# ================================================================
add_rect(slide, 0, 0, 45, 11.0, NAVY)

# Title
tb = slide.shapes.add_textbox(Cm(1.5), Cm(0.7), Cm(42), Cm(4.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = 'AI支援型コーディング技術（Vibe coding）による\n乳房形態評価アプリの開発'
run.font.size = Pt(29); run.font.bold = True
run.font.color.rgb = WHITE; run.font.name = FONT

# Authors
tb = slide.shapes.add_textbox(Cm(1.5), Cm(6.1), Cm(42), Cm(1.7))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '〇大槻祐喜、水野まどか、藤原久留望、浅香明紀、塗隆志'
run.font.size = Pt(18); run.font.color.rgb = WHITE; run.font.name = FONT

# Affiliation
tb = slide.shapes.add_textbox(Cm(1.5), Cm(7.9), Cm(42), Cm(1.4))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '大阪医科薬科大学　形成外科'
run.font.size = Pt(15); run.font.color.rgb = LT_BLUE; run.font.name = FONT

# Conference bar
add_rect(slide, 0, 11.0, 45, 1.4, MED_BLUE)
tb = slide.shapes.add_textbox(Cm(1), Cm(11.1), Cm(43), Cm(1.2))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '第69回日本形成外科学会総会学術集会　2026年4月22日（水）　一般演題「人工知能（AI)」　第9会場'
run.font.size = Pt(10.5); run.font.color.rgb = WHITE; run.font.name = FONT

# COI strip
add_rect(slide, 0, 12.4, 45, 0.5, LGRAY, BORDER, 0.3)
tb = slide.shapes.add_textbox(Cm(0.5), Cm(12.4), Cm(44), Cm(0.5))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
run = p.add_run()
run.text = '利益相反：なし'
run.font.size = Pt(8.5); run.font.color.rgb = RGBColor(0x55,0x55,0x55); run.font.name = FONT

# ================================================================
# ①  背景・目的   12.9 – 21.2 cm
# ================================================================
y = 12.9
sec_hdr(slide, 0.5, y, 44, '① 背景・目的', fs=14)
y += 0.9 + 0.2   # y = 14.0

# Slide 3 figure: original 12.3×11.5, AR=1.069  →  7.5×7.0 cm
FIG3_H, FIG3_W = 7.0, 7.0 * (12.3/11.5)   # 7.0 × 7.49
fig(slide, IMGS+'slide3_図 22.png', 1.0, y, w=FIG3_W, h=FIG3_H)

# Text alongside
add_lines(slide, 1.0+FIG3_W+0.5, y, 44.0-FIG3_W-2.0, FIG3_H, [
    ('【従来の課題】', True, 11.5),
    ('3D写真解析でも左右差の客観的定量化は困難', False, 10.5),
    ('測定が煩雑で臨床応用に乏しく', False, 10.5),
    ('差の臨床的意義も不明確', False, 10.5),
    ('', False, 4),
    ('【目　的】', True, 11.5),
    ('再建乳房の形態を　客観的に・簡単に', False, 11.0),
    ('評価できるソフトをAI（Vibe coding）で開発', True, 11.5),
])
# section bottom
y = 21.2

# ================================================================
# ②  Vibe coding によるアプリ開発   21.2 – 38.0 cm
# ================================================================
sec_hdr(slide, 0.5, y, 44, '② Vibe coding（AI支援型コーディング）によるアプリ開発', fs=13)
y += 0.9 + 0.3   # y = 22.4

# Slide 5: 23.3×14.1  AR=1.652  → height 14.5 cm, width 24.0 cm
FIG5_H = 14.5
FIG5_W = FIG5_H * (23.3/14.1)   # 24.0 cm
FIG5_L = (45.0 - FIG5_W) / 2
fig(slide, IMGS+'slide5_図 3.png', FIG5_L, y, w=FIG5_W, h=FIG5_H)
y += FIG5_H + 0.2   # y = 37.1

add_tb(slide, 1.5, y, 42, 0.7,
       'Cursor（Anysphere）でChatGPTに自然言語で指示→自動コード生成→解析ボタン1クリックで乳房左右差を算出',
       fs=10, align=PP_ALIGN.CENTER)
y = 38.0

# ================================================================
# ③  開発したアプリの画面   38.0 – 55.5 cm
# ================================================================
sec_hdr(slide, 0.5, y, 44, '③ 開発したアプリの画面', fs=14)
y += 0.9 + 0.3   # y = 39.2

# Slide 6: 21.1×14.5  AR=1.455  → height 15.0 cm, width 21.8 cm
FIG6_H = 15.0
FIG6_W = FIG6_H * (21.1/14.5)   # 21.8 cm
FIG6_L = (45.0 - FIG6_W) / 2
fig(slide, IMGS+'slide6_図 2.jpg', FIG6_L, y, w=FIG6_W, h=FIG6_H)
y += FIG6_H + 0.2   # y = 54.4

add_tb(slide, 1.5, y, 42, 0.7,
       '3Dデータ（STLファイル）の読み込み・乳房トリミング・解析を1クリックで実行するアプリ画面',
       fs=10, align=PP_ALIGN.CENTER)
y = 55.5

# ================================================================
# ④  BSAスコア算出アルゴリズム   55.5 – 72.7 cm
# ================================================================
sec_hdr(slide, 0.5, y, 44, '④ BSAスコア（Breast Symmetry Analyzer score）算出アルゴリズム', fs=13)
y += 0.9 + 0.3   # y = 56.7

# Slide 7 image: 22.5×14.3  AR=1.573  → height 14.5 cm, width 22.8 cm
FIG7_H = 14.5
FIG7_W = FIG7_H * (22.5/14.3)   # 22.8 cm
FIG7_L = (45.0 - FIG7_W) / 2
fig(slide, IMGS+'slide7_図 16.png', FIG7_L, y, w=FIG7_W, h=FIG7_H)
y += FIG7_H + 0.2   # y = 71.4

add_tb(slide, 1.5, y, 42, 1.0,
       '① STL入力・トリミング　→　② PCAでミラー平面生成・乳房重ね合わせ　→　③ KD-treeで頂点対応'
       '　→　④ RMS距離差を0〜100に正規化（BSAスコア）　→　⑤ 3Dカラーマップ可視化',
       fs=10, align=PP_ALIGN.CENTER)
y = 72.7

# ================================================================
# ⑤  精度評価・臨床結果   72.7 – 84.5 cm
# ================================================================
sec_hdr(slide, 0.5, y, 44, '⑤ 精度評価・臨床症例でのBSAスコア', fs=13)
y += 0.9 + 0.3   # y = 73.9

# ── Left half: Mannequin image + precision tables ──────────────
LW, LX = 21.5, 1.0    # left column
RW, RX = 21.5, 23.0   # right column

# Slide 8: 8.0×5.7  AR=1.404  →  height 6.0 cm, width 8.4 cm
FIG8_H, FIG8_W = 6.0, 6.0*(8.0/5.7)
FIG8_L = LX + (LW - FIG8_W) / 2
fig(slide, IMGS+'slide8_図 5.png', FIG8_L, y, w=FIG8_W, h=FIG8_H)

# Slide 11: 12.3×10.2  AR=1.206  →  height 6.0 cm, width 7.2 cm
FIG11_H, FIG11_W = 6.0, 6.0*(12.3/10.2)
FIG11_L = RX + (RW - FIG11_W) / 2
fig(slide, IMGS+'slide11_図 7.png', FIG11_L, y, w=FIG11_W, h=FIG11_H)

ty = y + FIG8_H + 0.3   # y = 80.2

# Left text: precision
add_lines(slide, LX, ty, LW, 3.8, [
    ('【精度評価：仮想乳房モデル】', True, 10.5),
    ('評価者A: BSAスコア 83.16±0.28（CV 0.34%）', False, 10),
    ('評価者B: BSAスコア 83.3 ±0.68（CV 0.81%）', False, 10),
    ('解析時間 A: 31.3秒 / B: 29.7秒', False, 10),
    ('→ 高い再現性・約30秒の迅速解析を確認', True, 10.5),
])

# Right text: clinical
add_lines(slide, RX, ty, RW, 3.8, [
    ('【臨床症例 n=81, 術後24ヶ月 VECTRA H2®】', True, 10.5),
    ('LD（n=33）: BSA mean=83.0±9.66', False, 10),
    ('TRAM（n=11）:BSA mean=79.7±8.70', False, 10),
    ('DIEP（n=37）:BSA mean=82.0±9.5', False, 10),
    ('one-way ANOVA  F(2,78)=0.51, p=0.61（n.s.）', True, 10),
])

y = 84.5

# ================================================================
# ⑥  考察   84.5 – 91.0 cm
# ================================================================
sec_hdr(slide, 0.5, y, 44, '⑥ 考察', fs=14)
y += 0.9 + 0.25   # y = 85.65

add_lines(slide, 1.0, y, 43.0, 5.1, [
    ('【精度の誤差要因】', True, 11.0),
    ('  乳房ROI設定の主観性・手動トリミング誤差　→　機械学習による自動認識で解消可能', False, 10.5),
    ('', False, 3),
    ('【今後の展望】', True, 11.0),
    ('  ・主観的評価や患者満足度との相関確認', False, 10.5),
    ('  ・術式・放射線照射の有無など再建結果に影響する要因の解析', False, 10.5),
    ('  ・客観的数値指標として再建乳房評価を標準化', False, 10.5),
])

# ================================================================
# FOOTER  まとめ   91.0 – 95.0 cm
# ================================================================
add_rect(slide, 0, 91.0, 45, 4.0, MED_BLUE)

tb = slide.shapes.add_textbox(Cm(0.5), Cm(91.2), Cm(44), Cm(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = '■ まとめ'
run.font.size = Pt(15); run.font.bold = True
run.font.color.rgb = WHITE; run.font.name = FONT

tb = slide.shapes.add_textbox(Cm(1.0), Cm(92.3), Cm(43), Cm(2.5))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate([
    '① 乳房形態の定量評価により、症例間・術式間比較および経時的変化の定量化が可能となる',
    '② 1クリックでスコア化することで定量化の煩雑さを解決',
    '③ AIの発展によりさらに簡便で精度の高いスコアの開発が期待できる',
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(10.5); run.font.color.rgb = WHITE; run.font.name = FONT

# ================================================================
prs.save('poster_Vibe_coding_v2.pptx')
print(f'Saved: poster_Vibe_coding_v2.pptx  ({prs.slide_width.cm:.0f}x{prs.slide_height.cm:.0f} cm)')
