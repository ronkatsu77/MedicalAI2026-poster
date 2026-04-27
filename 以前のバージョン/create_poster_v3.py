import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

NAVY     = RGBColor(0x1F, 0x38, 0x64)
MED_BLUE = RGBColor(0x2E, 0x74, 0xB5)
LT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
FONT     = 'メイリオ'

# ════════════════════════════════════════════════════════
# 1. TITLE CARD  70 × 20 cm
# ════════════════════════════════════════════════════════
prs_t = Presentation()
prs_t.slide_width  = Cm(70)
prs_t.slide_height = Cm(20)

slide_t = prs_t.slides.add_slide(prs_t.slide_layouts[6])
slide_t.background.fill.solid()
slide_t.background.fill.fore_color.rgb = NAVY

def tb(slide, l, t, w, h, text, fs, bold=False, color=WHITE,
       align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(fs)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = FONT

# Title
tb(slide_t, 3, 1.5, 64, 7,
   'Vibe codingによる乳房形態評価アプリ開発の試み', 42, bold=True)

# Horizontal accent line
acc = slide_t.shapes.add_shape(1, Cm(12), Cm(9.8), Cm(46), Cm(0.15))
acc.fill.solid(); acc.fill.fore_color.rgb = LT_BLUE
acc.line.fill.background()

# Authors
tb(slide_t, 3, 11.0, 64, 3,
   '〇大槻祐喜、水野まどか、藤原久留望、浅香明紀、塗隆志', 26)

# Affiliation
tb(slide_t, 3, 15.0, 64, 3,
   '大阪医科薬科大学　形成外科', 22, color=LT_BLUE)

prs_t.save('3poster_title_70x20.pptx')
print('Saved: 3poster_title_70x20.pptx  (70x20cm)')


# ════════════════════════════════════════════════════════
# 2. MAIN POSTER  – ヘッダー削除 & コンテンツ再配置
# ════════════════════════════════════════════════════════
prs = Presentation('2poster_Vibe_coding_v2.pptx')
slide = prs.slides[0]

SLIDE_H    = prs.slide_height.cm        # 95 cm
SHIFT      = 11.5                       # header height to remove
DEL_THRESH = 11.0                       # shapes with top < this → delete
HEADER_KW  = ['利益相反', '第69回', '〇大槻', '大阪医科薬科大学　形成外科',
              '人工知能']

def is_header(s):
    """Return True if this shape belongs to the header block."""
    # ── NEVER delete content section markers ──────────────────
    if hasattr(s, 'text') and s.text.strip():
        t = s.text.strip()
        # Section headers (①②③…) and まとめ are always content
        if any(m in t for m in ['①', '②', '③', '④', '⑤', '⑥', '■']):
            return False

    # Shapes clearly inside the header zone (above content start)
    if s.top.cm < DEL_THRESH:
        return True

    # Conference bar / COI bar at DEL_THRESH – SHIFT+1.5 cm:
    # large-width rect (conference bar) or thin strip (COI bar) or keyword text
    if DEL_THRESH <= s.top.cm < SHIFT + 1.5:
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if s.width.cm > 40:          # large bg rect = conference bar
                return True
            if s.height.cm < 0.8:        # thin strip = COI bar
                return True
        if hasattr(s, 'text'):
            for kw in HEADER_KW:
                if kw in s.text:
                    return True
    return False

sp_tree = slide.shapes._spTree
to_del  = [s._element for s in slide.shapes if is_header(s)]
print(f'Deleting {len(to_del)} header shapes ...')
for sp in to_del:
    sp_tree.remove(sp)

# ── Shift all remaining shapes up by SHIFT cm ─────────────────
remaining = list(slide.shapes)
for s in remaining:
    s.top = Cm(s.top.cm - SHIFT)

# ── Find まとめ footer background rect & extend to bottom ─────
# Criteria: AUTO_SHAPE, width > 40cm, top > 75cm (after shift ~79.4cm),
# and EMPTY text (it's a pure background rect, no text of its own).
footer_rect = None
for s in remaining:
    if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.width.cm > 40:
        text = s.text.strip() if hasattr(s, 'text') else ''
        if s.top.cm > 75 and text == '':   # ← empty text = background rect
            footer_rect = s
            print(f'Found footer rect at top={s.top.cm:.1f}cm (h={s.height.cm:.1f}cm)')
            break

if footer_rect:
    f_top = footer_rect.top.cm
    new_h = SLIDE_H - f_top
    footer_rect.height = Cm(new_h)
    print(f'Footer extended: {f_top:.1f} → {SLIDE_H:.0f}cm  (height {new_h:.1f}cm)')

    # Reposition まとめ content within the taller footer:
    #   "■ まとめ" title → f_top + 2.0 cm
    #   summary text    → f_top + 4.5 cm  (with height to reach slide bottom)
    # Only touch shapes that are currently inside the original footer region
    # (between f_top−0.5 and f_top+5.0 cm, excluding the bg rect itself).
    footer_content = [s for s in remaining
                      if s is not footer_rect
                      and f_top - 0.5 <= s.top.cm <= f_top + 5.0]
    footer_content.sort(key=lambda s: s.top.cm)
    print(f'まとめ内テキスト shapes: {len(footer_content)}')
    for s in footer_content:
        text = s.text.strip() if hasattr(s, 'text') else ''
        if '■' in text and 'まとめ' in text:
            s.top = Cm(f_top + 2.0)
        else:
            s.top    = Cm(f_top + 4.5)
            s.height = Cm(SLIDE_H - (f_top + 4.5) - 1.0)

print('Shapes after edit:', len(list(slide.shapes)))
prs.save('3poster_main.pptx')
print(f'Saved: 3poster_main.pptx  ({prs.slide_width.cm:.0f}x{SLIDE_H:.0f}cm)')
