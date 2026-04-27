import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('3poster_main.pptx')
slide = prs.slides[0]
SLIDE_H = prs.slide_height.cm  # 95.0

# ── Old section boundaries (top of each section header) ──────────────
# Measured from create_poster_v3.py output
OLD_BOUNDS = [0.04, 8.62, 19.57, 49.92, 72.26, 79.40, 95.0]
# Sections: ①[0] ②[1] ③[2] ④[3] ⑤[4] まとめ[5]

# ── New layout: shrink まとめ to 9 cm, scale ①–⑤ proportionally ─────
MATOME_H_NEW = 9.0
old_content = OLD_BOUNDS[5] - OLD_BOUNDS[0]  # 79.36 cm
new_content  = SLIDE_H - MATOME_H_NEW         # 86.0 cm
SCALE = new_content / old_content             # ≈1.083

NEW_BOUNDS = [0.0]
for i in range(5):
    old_h = OLD_BOUNDS[i + 1] - OLD_BOUNDS[i]
    NEW_BOUNDS.append(NEW_BOUNDS[-1] + old_h * SCALE)
NEW_BOUNDS.append(SLIDE_H)

for i, (nb, ob) in enumerate(zip(NEW_BOUNDS, OLD_BOUNDS)):
    labels = ['①','②','③','④','⑤','まとめ','END']
    print(f'  {labels[i]:4s}  old={ob:.2f}  new={nb:.2f}')


def get_sec(top):
    """Return section index (0–5) for a shape whose top is `top` cm."""
    EPS = 0.5
    for i in range(6):
        if OLD_BOUNDS[i] - EPS <= top < OLD_BOUNDS[i + 1]:
            return i
    # Anything at or past the last boundary belongs to まとめ
    return 5


def process(s):
    old_top = s.top.cm
    old_h   = s.height.cm
    old_w   = s.width.cm
    stype   = s.shape_type

    i       = get_sec(old_top)
    old_s   = OLD_BOUNDS[i]
    old_e   = OLD_BOUNDS[i + 1]
    new_s   = NEW_BOUNDS[i]
    new_e   = NEW_BOUNDS[i + 1]
    sec_s   = (new_e - new_s) / (old_e - old_s)  # per-section scale factor

    # Move top
    new_top = new_s + (old_top - old_s) * sec_s
    s.top   = Cm(max(0.0, new_top))

    # Scale height / width based on shape type
    if stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
        # Wide thin bars = section header strip → keep height
        if old_w > 40 and old_h < 2.0:
            return
        s.height = Cm(old_h * sec_s)

    elif stype == MSO_SHAPE_TYPE.TEXT_BOX:
        # Only scale taller text boxes (section content); keep small labels as-is
        if old_h >= 2.0:
            s.height = Cm(old_h * sec_s)

    elif stype == MSO_SHAPE_TYPE.PICTURE:
        s.height = Cm(old_h * sec_s)
        # Scale width proportionally to preserve aspect ratio
        s.width  = Cm(old_w * sec_s)

    elif stype == MSO_SHAPE_TYPE.GROUP:
        s.height = Cm(old_h * sec_s)
        s.width  = Cm(old_w * sec_s)

    elif stype == MSO_SHAPE_TYPE.LINE:
        if old_h > 0.5:
            s.height = Cm(old_h * sec_s)

    elif stype == MSO_SHAPE_TYPE.TABLE:
        new_h_cm = old_h * sec_s
        s.height = Cm(new_h_cm)
        tbl = s.table
        total = sum(r.height for r in tbl.rows)
        if total > 0:
            new_h_emu = s.height
            for r in tbl.rows:
                r.height = int(r.height / total * new_h_emu)


for shape in slide.shapes:
    process(shape)

print(f'\nShapes processed: {len(list(slide.shapes))}')
prs.save('3poster_v4.pptx')
print('Saved: 3poster_v4.pptx')
